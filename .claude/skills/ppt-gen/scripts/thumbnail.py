#!/usr/bin/env python3
"""
Create thumbnail grids from PowerPoint presentation slides.

Creates a grid layout of slide thumbnails with configurable columns (max 6).
Each grid contains up to cols×(cols+1) images. For presentations with more
slides, multiple numbered grid files are created automatically.

The program outputs the names of all files created.

Output:
- Single grid: {prefix}.jpg (if slides fit in one grid)
- Multiple grids: {prefix}-1.jpg, {prefix}-2.jpg, etc.

Grid limits by column count:
- 3 cols: max 12 slides per grid (3×4)
- 4 cols: max 20 slides per grid (4×5)
- 5 cols: max 30 slides per grid (5×6) [default]
- 6 cols: max 42 slides per grid (6×7)

Usage:
    python thumbnail.py input.pptx [output_prefix] [--cols N] [--outline-placeholders]
    python thumbnail.py input.pptx output/ --slides 5 --single

Examples:
    python thumbnail.py presentation.pptx
    # Creates: thumbnails.jpg (using default prefix)
    # Outputs:
    #   Created 1 grid(s):
    #     - thumbnails.jpg

    python thumbnail.py large-deck.pptx grid --cols 4
    # Creates: grid-1.jpg, grid-2.jpg, grid-3.jpg
    # Outputs:
    #   Created 3 grid(s):
    #     - grid-1.jpg
    #     - grid-2.jpg
    #     - grid-3.jpg

    python thumbnail.py template.pptx analysis --outline-placeholders
    # Creates thumbnail grids with red outlines around text placeholders

    python thumbnail.py input.pptx output/ --slides 5 --single
    # Creates: output/slide-5.png (single slide at 1980x1080)
    # Outputs:
    #   Created single thumbnail:
    #     - output/slide-5.png

    python thumbnail.py input.pptx output/ --slides 1,3,5
    # Creates: output/slide-1.png, output/slide-3.png, output/slide-5.png
"""

# Single slide output constants
SINGLE_WIDTH = 1980
SINGLE_HEIGHT = 1080

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path

from inventory import extract_text_inventory
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

# Constants
THUMBNAIL_WIDTH = 300  # Fixed thumbnail width in pixels
CONVERSION_DPI = 100  # DPI for PDF to image conversion
MAX_COLS = 6  # Maximum number of columns
DEFAULT_COLS = 5  # Default number of columns
JPEG_QUALITY = 95  # JPEG compression quality

# Grid layout constants
GRID_PADDING = 20  # Padding between thumbnails
BORDER_WIDTH = 2  # Border width around thumbnails
FONT_SIZE_RATIO = 0.12  # Font size as fraction of thumbnail width
LABEL_PADDING_RATIO = 0.4  # Label padding as fraction of font size


def main():
    parser = argparse.ArgumentParser(
        description="Create thumbnail grids from PowerPoint slides."
    )
    parser.add_argument("input", help="Input PowerPoint file (.pptx)")
    parser.add_argument(
        "output_prefix",
        nargs="?",
        default="thumbnails",
        help="Output prefix for image files (default: thumbnails, will create prefix.jpg or prefix-N.jpg)",
    )
    parser.add_argument(
        "--cols",
        type=int,
        default=DEFAULT_COLS,
        help=f"Number of columns (default: {DEFAULT_COLS}, max: {MAX_COLS})",
    )
    parser.add_argument(
        "--outline-placeholders",
        action="store_true",
        help="Outline text placeholders with a colored border",
    )
    parser.add_argument(
        "--slides",
        type=str,
        help="Specific slide(s) to extract (0-based index). Use comma for multiple: --slides 5 or --slides 1,3,5",
    )
    parser.add_argument(
        "--single",
        action="store_true",
        help="Output single slide image(s) instead of grid (1980x1080 PNG)",
    )

    args = parser.parse_args()

    # Validate columns
    cols = min(args.cols, MAX_COLS)
    if args.cols > MAX_COLS:
        print(f"Warning: Columns limited to {MAX_COLS} (requested {args.cols})")

    # Validate input
    input_path = Path(args.input)
    if not input_path.exists() or input_path.suffix.lower() != ".pptx":
        print(f"Error: Invalid PowerPoint file: {args.input}")
        sys.exit(1)

    # Parse slide indices if specified
    slide_indices = None
    if args.slides:
        try:
            slide_indices = [int(s.strip()) for s in args.slides.split(",")]
        except ValueError:
            print(f"Error: Invalid slide indices: {args.slides}")
            sys.exit(1)

    # Handle single slide mode
    if args.single or slide_indices:
        output_dir = Path(args.output_prefix)
        output_dir.mkdir(parents=True, exist_ok=True)
        create_single_thumbnails(input_path, output_dir, slide_indices, args.single)
        return

    # Construct output path (always JPG for grid mode)
    output_path = Path(f"{args.output_prefix}.jpg")

    print(f"Processing: {args.input}")

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            # Get placeholder regions if outlining is enabled
            placeholder_regions = None
            slide_dimensions = None
            if args.outline_placeholders:
                print("Extracting placeholder regions...")
                placeholder_regions, slide_dimensions = get_placeholder_regions(
                    input_path
                )
                if placeholder_regions:
                    print(f"Found placeholders on {len(placeholder_regions)} slides")

            # Convert slides to images
            slide_images = convert_to_images(input_path, Path(temp_dir), CONVERSION_DPI)
            if not slide_images:
                print("Error: No slides found")
                sys.exit(1)

            print(f"Found {len(slide_images)} slides")

            # Create grids (max cols×(cols+1) images per grid)
            grid_files = create_grids(
                slide_images,
                cols,
                THUMBNAIL_WIDTH,
                output_path,
                placeholder_regions,
                slide_dimensions,
            )

            # Print saved files
            print(f"Created {len(grid_files)} grid(s):")
            for grid_file in grid_files:
                print(f"  - {grid_file}")

    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


def create_hidden_slide_placeholder(size):
    """Create placeholder image for hidden slides."""
    img = Image.new("RGB", size, color="#F0F0F0")
    draw = ImageDraw.Draw(img)
    line_width = max(5, min(size) // 100)
    draw.line([(0, 0), size], fill="#CCCCCC", width=line_width)
    draw.line([(size[0], 0), (0, size[1])], fill="#CCCCCC", width=line_width)
    return img


def create_single_thumbnails(pptx_path, output_dir, slide_indices=None, single_mode=False):
    """Create single slide thumbnail(s) at high resolution.

    Args:
        pptx_path: Path to the PowerPoint file
        output_dir: Output directory for thumbnails
        slide_indices: List of slide indices (0-based) to extract. If None, extracts all.
        single_mode: If True and slide_indices has one item, outputs single file
    """
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir = Path(temp_dir)

        # Get total slide count and hidden slides
        prs = Presentation(str(pptx_path))
        total_slides = len(prs.slides)
        hidden_slides = {
            idx for idx, slide in enumerate(prs.slides)
            if slide.element.get("show") == "0"
        }

        # Default to all slides if not specified
        if slide_indices is None:
            slide_indices = list(range(total_slides))

        # Validate indices
        valid_indices = [i for i in slide_indices if 0 <= i < total_slides]
        if not valid_indices:
            print(f"Error: No valid slide indices (total slides: {total_slides})")
            sys.exit(1)

        # Convert to PDF
        pdf_path = temp_dir / f"{pptx_path.stem}.pdf"
        print(f"Converting to PDF...")
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(temp_dir),
                str(pptx_path),
            ],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0 or not pdf_path.exists():
            raise RuntimeError("PDF conversion failed")

        # Calculate DPI for target resolution (1980x1080)
        # Standard slide is 10" x 7.5" (16:9 aspect)
        # DPI = pixels / inches → 1980 / 10 = 198 DPI
        target_dpi = 198

        # Convert PDF to images at high resolution
        print(f"Converting to images at {target_dpi} DPI...")
        result = subprocess.run(
            ["pdftoppm", "-png", "-r", str(target_dpi), str(pdf_path), str(temp_dir / "slide")],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            raise RuntimeError("Image conversion failed")

        # Get visible slide images (pdftoppm skips hidden slides)
        visible_images = sorted(temp_dir.glob("slide-*.png"))

        # Map visible image index to actual slide index
        visible_to_actual = {}
        visible_idx = 0
        for actual_idx in range(total_slides):
            if actual_idx not in hidden_slides:
                visible_to_actual[actual_idx] = visible_idx
                visible_idx += 1

        # Process requested slides
        created_files = []
        for slide_idx in valid_indices:
            output_path = output_dir / f"slide-{slide_idx}.png"

            if slide_idx in hidden_slides:
                # Create placeholder for hidden slide
                placeholder_img = create_hidden_slide_placeholder((SINGLE_WIDTH, SINGLE_HEIGHT))
                placeholder_img.save(output_path, "PNG")
                created_files.append(str(output_path))
                print(f"  Created placeholder for hidden slide {slide_idx}")
            else:
                # Get the corresponding visible image
                vis_idx = visible_to_actual.get(slide_idx)
                if vis_idx is not None and vis_idx < len(visible_images):
                    with Image.open(visible_images[vis_idx]) as img:
                        # Resize to target dimensions
                        img_resized = img.resize(
                            (SINGLE_WIDTH, SINGLE_HEIGHT),
                            Image.Resampling.LANCZOS
                        )
                        img_resized.save(output_path, "PNG")
                        created_files.append(str(output_path))
                        print(f"  Created: {output_path}")
                else:
                    print(f"  Warning: Could not find image for slide {slide_idx}")

        # Print summary
        if len(created_files) == 1:
            print(f"Created single thumbnail:")
        else:
            print(f"Created {len(created_files)} thumbnail(s):")
        for f in created_files:
            print(f"  - {f}")


def get_placeholder_regions(pptx_path):
    """Extract ALL text regions from the presentation.

    Returns a tuple of (placeholder_regions, slide_dimensions).
    text_regions is a dict mapping slide indices to lists of text regions.
    Each region is a dict with 'left', 'top', 'width', 'height' in inches.
    slide_dimensions is a tuple of (width_inches, height_inches).
    """
    prs = Presentation(str(pptx_path))
    inventory = extract_text_inventory(pptx_path, prs)
    placeholder_regions = {}

    # Get actual slide dimensions in inches (EMU to inches conversion)
    slide_width_inches = (prs.slide_width or 9144000) / 914400.0
    slide_height_inches = (prs.slide_height or 5143500) / 914400.0

    for slide_key, shapes in inventory.items():
        # Extract slide index from "slide-N" format
        slide_idx = int(slide_key.split("-")[1])
        regions = []

        for shape_key, shape_data in shapes.items():
            # The inventory only contains shapes with text, so all shapes should be highlighted
            regions.append(
                {
                    "left": shape_data.left,
                    "top": shape_data.top,
                    "width": shape_data.width,
                    "height": shape_data.height,
                }
            )

        if regions:
            placeholder_regions[slide_idx] = regions

    return placeholder_regions, (slide_width_inches, slide_height_inches)


def convert_to_images(pptx_path, temp_dir, dpi):
    """Convert PowerPoint to images via PDF, handling hidden slides."""
    # Detect hidden slides
    print("Analyzing presentation...")
    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)

    # Find hidden slides (1-based indexing for display)
    hidden_slides = {
        idx + 1
        for idx, slide in enumerate(prs.slides)
        if slide.element.get("show") == "0"
    }

    print(f"Total slides: {total_slides}")
    if hidden_slides:
        print(f"Hidden slides: {sorted(hidden_slides)}")

    pdf_path = temp_dir / f"{pptx_path.stem}.pdf"

    # Convert to PDF
    print("Converting to PDF...")
    result = subprocess.run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(temp_dir),
            str(pptx_path),
        ],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0 or not pdf_path.exists():
        raise RuntimeError("PDF conversion failed")

    # Convert PDF to images
    print(f"Converting to images at {dpi} DPI...")
    result = subprocess.run(
        ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf_path), str(temp_dir / "slide")],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        raise RuntimeError("Image conversion failed")

    visible_images = sorted(temp_dir.glob("slide-*.jpg"))

    # Create full list with placeholders for hidden slides
    all_images = []
    visible_idx = 0

    # Get placeholder dimensions from first visible slide
    if visible_images:
        with Image.open(visible_images[0]) as img:
            placeholder_size = img.size
    else:
        placeholder_size = (1920, 1080)

    for slide_num in range(1, total_slides + 1):
        if slide_num in hidden_slides:
            # Create placeholder image for hidden slide
            placeholder_path = temp_dir / f"hidden-{slide_num:03d}.jpg"
            placeholder_img = create_hidden_slide_placeholder(placeholder_size)
            placeholder_img.save(placeholder_path, "JPEG")
            all_images.append(placeholder_path)
        else:
            # Use the actual visible slide image
            if visible_idx < len(visible_images):
                all_images.append(visible_images[visible_idx])
                visible_idx += 1

    return all_images


def create_grids(
    image_paths,
    cols,
    width,
    output_path,
    placeholder_regions=None,
    slide_dimensions=None,
):
    """Create multiple thumbnail grids from slide images, max cols×(cols+1) images per grid."""
    # Maximum images per grid is cols × (cols + 1) for better proportions
    max_images_per_grid = cols * (cols + 1)
    grid_files = []

    print(
        f"Creating grids with {cols} columns (max {max_images_per_grid} images per grid)"
    )

    # Split images into chunks
    for chunk_idx, start_idx in enumerate(
        range(0, len(image_paths), max_images_per_grid)
    ):
        end_idx = min(start_idx + max_images_per_grid, len(image_paths))
        chunk_images = image_paths[start_idx:end_idx]

        # Create grid for this chunk
        grid = create_grid(
            chunk_images, cols, width, start_idx, placeholder_regions, slide_dimensions
        )

        # Generate output filename
        if len(image_paths) <= max_images_per_grid:
            # Single grid - use base filename without suffix
            grid_filename = output_path
        else:
            # Multiple grids - insert index before extension with dash
            stem = output_path.stem
            suffix = output_path.suffix
            grid_filename = output_path.parent / f"{stem}-{chunk_idx + 1}{suffix}"

        # Save grid
        grid_filename.parent.mkdir(parents=True, exist_ok=True)
        grid.save(str(grid_filename), quality=JPEG_QUALITY)
        grid_files.append(str(grid_filename))

    return grid_files


def create_grid(
    image_paths,
    cols,
    width,
    start_slide_num=0,
    placeholder_regions=None,
    slide_dimensions=None,
):
    """Create thumbnail grid from slide images with optional placeholder outlining."""
    font_size = int(width * FONT_SIZE_RATIO)
    label_padding = int(font_size * LABEL_PADDING_RATIO)

    # Get dimensions
    with Image.open(image_paths[0]) as img:
        aspect = img.height / img.width
    height = int(width * aspect)

    # Calculate grid size
    rows = (len(image_paths) + cols - 1) // cols
    grid_w = cols * width + (cols + 1) * GRID_PADDING
    grid_h = rows * (height + font_size + label_padding * 2) + (rows + 1) * GRID_PADDING

    # Create grid
    grid = Image.new("RGB", (grid_w, grid_h), "white")
    draw = ImageDraw.Draw(grid)

    # Load font with size based on thumbnail width
    try:
        # Use Pillow's default font with size
        font = ImageFont.load_default(size=font_size)
    except Exception:
        # Fall back to basic default font if size parameter not supported
        font = ImageFont.load_default()

    # Place thumbnails
    for i, img_path in enumerate(image_paths):
        row, col = i // cols, i % cols
        x = col * width + (col + 1) * GRID_PADDING
        y_base = (
            row * (height + font_size + label_padding * 2) + (row + 1) * GRID_PADDING
        )

        # Add label with actual slide number
        label = f"{start_slide_num + i}"
        bbox = draw.textbbox((0, 0), label, font=font)
        text_w = bbox[2] - bbox[0]
        draw.text(
            (x + (width - text_w) // 2, y_base + label_padding),
            label,
            fill="black",
            font=font,
        )

        # Add thumbnail below label with proportional spacing
        y_thumbnail = y_base + label_padding + font_size + label_padding

        with Image.open(img_path) as img:
            # Get original dimensions before thumbnail
            orig_w, orig_h = img.size

            # Apply placeholder outlines if enabled
            if placeholder_regions and (start_slide_num + i) in placeholder_regions:
                # Convert to RGBA for transparency support
                if img.mode != "RGBA":
                    img = img.convert("RGBA")

                # Get the regions for this slide
                regions = placeholder_regions[start_slide_num + i]

                # Calculate scale factors using actual slide dimensions
                if slide_dimensions:
                    slide_width_inches, slide_height_inches = slide_dimensions
                else:
                    # Fallback: estimate from image size at CONVERSION_DPI
                    slide_width_inches = orig_w / CONVERSION_DPI
                    slide_height_inches = orig_h / CONVERSION_DPI

                x_scale = orig_w / slide_width_inches
                y_scale = orig_h / slide_height_inches

                # Create a highlight overlay
                overlay = Image.new("RGBA", img.size, (255, 255, 255, 0))
                overlay_draw = ImageDraw.Draw(overlay)

                # Highlight each placeholder region
                for region in regions:
                    # Convert from inches to pixels in the original image
                    px_left = int(region["left"] * x_scale)
                    px_top = int(region["top"] * y_scale)
                    px_width = int(region["width"] * x_scale)
                    px_height = int(region["height"] * y_scale)

                    # Draw highlight outline with red color and thick stroke
                    # Using a bright red outline instead of fill
                    stroke_width = max(
                        5, min(orig_w, orig_h) // 150
                    )  # Thicker proportional stroke width
                    overlay_draw.rectangle(
                        [(px_left, px_top), (px_left + px_width, px_top + px_height)],
                        outline=(255, 0, 0, 255),  # Bright red, fully opaque
                        width=stroke_width,
                    )

                # Composite the overlay onto the image using alpha blending
                img = Image.alpha_composite(img, overlay)
                # Convert back to RGB for JPEG saving
                img = img.convert("RGB")

            img.thumbnail((width, height), Image.Resampling.LANCZOS)
            w, h = img.size
            tx = x + (width - w) // 2
            ty = y_thumbnail + (height - h) // 2
            grid.paste(img, (tx, ty))

            # Add border
            if BORDER_WIDTH > 0:
                draw.rectangle(
                    [
                        (tx - BORDER_WIDTH, ty - BORDER_WIDTH),
                        (tx + w + BORDER_WIDTH - 1, ty + h + BORDER_WIDTH - 1),
                    ],
                    outline="gray",
                    width=BORDER_WIDTH,
                )

    return grid


if __name__ == "__main__":
    main()
